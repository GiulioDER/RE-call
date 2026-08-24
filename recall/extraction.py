"""Format aware document extraction for the desktop and MCP ingest paths.

The index stores UTF 8 text, but users should not have to convert a document before adding it.
This module keeps the original file as the manifest object and derives a searchable text view for
UTF 8 text, PDF, DOCX, XLSX, PPTX, EML, MSG, RTF, EPUB, HTML, CSV, TSV, and legacy Office files.
DOC, ODT, ODS, ODP, and PPT require LibreOffice. MSG requires the optional documents extra.
"""

from __future__ import annotations

import atexit
import contextlib
import csv
import io
import os
import posixpath
import queue
import re
import shutil
import subprocess
import tempfile
import threading
import urllib.parse
import zipfile
from dataclasses import dataclass
from email import policy
from email.parser import BytesParser
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Iterable, Iterator, Literal
from xml.etree import ElementTree
from recall.errors import RecallError


class DocumentExtractionError(ValueError, RecallError):
    """Raised when a supported file cannot produce a searchable text view."""


@dataclass(frozen=True)
class ExtractedBlock:
    """A typed searchable region derived from one source object.

    ``kind`` is ``text`` or ``table``. Table metadata contains ``content_kind=table``,
    ``table_headers``, and, when numeric cells are present, ``numeric_values``. The block shape
    is versioned by ``STRUCTURED_DOCUMENT_VERSION`` because it is persisted in chunk metadata.
    """

    text: str
    kind: Literal["text", "table"]
    metadata: dict[str, Any]


@dataclass(frozen=True)
class ExtractedDocument:
    text: str
    media_type: str
    metadata: dict[str, Any]
    blocks: tuple[ExtractedBlock, ...] = ()


TEXT_EXTENSIONS = frozenset(
    {
        ".c",
        ".cc",
        ".cpp",
        ".cs",
        ".css",
        ".csv",
        ".go",
        ".h",
        ".hpp",
        ".html",
        ".htm",
        ".java",
        ".js",
        ".jsx",
        ".json",
        ".md",
        ".markdown",
        ".mdx",
        ".py",
        ".pyi",
        ".rb",
        ".rs",
        ".rst",
        ".sh",
        ".sql",
        ".svelte",
        ".toml",
        ".ts",
        ".tsx",
        ".txt",
        ".tsv",
        ".xml",
        ".yaml",
        ".yml",
    }
)
DOCUMENT_EXTENSIONS = frozenset(
    {
        *TEXT_EXTENSIONS,
        ".doc",
        ".docm",
        ".docx",
        ".dotm",
        ".dotx",
        ".eml",
        ".epub",
        ".msg",
        ".odt",
        ".odp",
        ".ods",
        ".pdf",
        ".ppt",
        ".pptm",
        ".pptx",
        ".potm",
        ".potx",
        ".rtf",
        ".xls",
        ".xlsm",
        ".xlsx",
        ".xltm",
        ".xltx",
    }
)

# The formats handed to the LibreOffice CLI because no Python reader covers them. MSG is
# deliberately absent: LibreOffice ships no MAPI import filter, so a .msg reaches it only to fail.
# Measured 2026-08-18 against LibreOffice 25.8, exit 1 and "source file could not be loaded" on a
# genuine .msg; re-measure with scripts/check_libreoffice_msg.py. MSG is read by python oxmsg from
# the documents extra, and a deployment without it is told to install that extra rather than being
# routed here and told to install LibreOffice.
LIBREOFFICE_EXTENSIONS = frozenset({".doc", ".odt", ".ods", ".odp", ".ppt"})

MAX_TABLE_ROWS = 10_000
MAX_TABLE_COLUMNS = 200
MAX_TABLE_CELLS = 500_000
MAX_EXTRACTED_CHARACTERS = 5_000_000
MAX_EPUB_MEMBERS = 10_000
MAX_EPUB_MEMBER_BYTES = MAX_EXTRACTED_CHARACTERS * 4
MAX_EPUB_TOTAL_BYTES = MAX_EXTRACTED_CHARACTERS * 4
STRUCTURED_DOCUMENT_VERSION = "table-row-groups-v1"

_MEDIA_TYPE_SUFFIXES = {
    "text/plain": ".txt",
    "text/markdown": ".md",
    "text/csv": ".csv",
    "text/tab-separated-values": ".tsv",
    "application/json": ".json",
    "application/xml": ".xml",
    "text/xml": ".xml",
    "application/pdf": ".pdf",
}


def extraction_supported(path: str | Path) -> bool:
    return Path(path).suffix.lower() in DOCUMENT_EXTENSIONS


def extraction_path_for(path: str | Path, media_type: str) -> Path:
    candidate = Path(path)
    if candidate.suffix.lower() in DOCUMENT_EXTENSIONS:
        return candidate
    suffix = _MEDIA_TYPE_SUFFIXES.get(media_type.lower().split(";", 1)[0].strip())
    return candidate.with_name(candidate.name + suffix) if suffix else candidate


def extract_document(path: Path, data: bytes) -> ExtractedDocument:
    """Return a UTF 8 searchable representation while retaining source provenance."""
    suffix = path.suffix.lower()
    if suffix not in DOCUMENT_EXTENSIONS:
        raise DocumentExtractionError(
            f"unsupported file type {suffix or '<no extension>'}; supported document formats "
            "include PDF, DOCX, XLSX, PPTX, CSV, HTML, and UTF 8 text"
        )

    if suffix == ".pdf":
        return _extract_pdf(path, data)
    if suffix in {".docx", ".docm", ".dotx", ".dotm"}:
        return _extract_docx(data)
    if suffix in {".xlsx", ".xlsm", ".xltx", ".xltm", ".xls"}:
        return _extract_spreadsheet(suffix, data)
    if suffix in {".pptx", ".pptm", ".potx", ".potm"}:
        return _extract_pptx(data)
    if suffix == ".eml":
        return _extract_email(data)
    if suffix == ".msg":
        return _extract_msg(data)
    if suffix == ".rtf":
        return _extract_rtf(data)
    if suffix == ".epub":
        return _extract_epub(data)
    if suffix in {".html", ".htm"}:
        return _extract_html(data)
    if suffix in LIBREOFFICE_EXTENSIONS:
        return _extract_with_libreoffice(path, data)
    if suffix in {".csv", ".tsv"}:
        return _extract_delimited(suffix, data)
    text = _decode_text(data)
    return _document(
        text,
        "text/markdown" if suffix in {".md", ".markdown", ".mdx"} else "text/plain",
        {"source_format": suffix.removeprefix("."), "extraction": "text"},
    )


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return data.decode(encoding).replace("\x00", "")
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace").replace("\x00", "")


def _result(text: str, source_format: str, *, tables: int = 0, **metadata: Any) -> ExtractedDocument:
    clean = text.strip()
    if not clean:
        raise DocumentExtractionError(
            f"{source_format.upper()} did not contain extractable text or tables"
        )
    if len(clean) > MAX_EXTRACTED_CHARACTERS:
        raise DocumentExtractionError(
            f"{source_format.upper()} extracted more than {MAX_EXTRACTED_CHARACTERS:,} characters; "
            "split the source into smaller files"
        )
    return _document(
        clean,
        "text/plain",
        {
            "source_format": source_format,
            "extraction": "structured" if tables else "text",
            "table_count": tables,
            **metadata,
        },
    )


def _document(text: str, media_type: str, metadata: dict[str, Any]) -> ExtractedDocument:
    clean = text.strip()
    source_format = str(metadata.get("source_format", "document"))
    if not clean:
        raise DocumentExtractionError(
            f"{source_format.upper()} did not contain extractable text or tables"
        )
    if len(clean) > MAX_EXTRACTED_CHARACTERS:
        raise DocumentExtractionError(
            f"{source_format.upper()} extracted more than {MAX_EXTRACTED_CHARACTERS:,} characters; "
            "split the source into smaller files"
        )
    return ExtractedDocument(clean, media_type, metadata, _blocks_from_text(clean, metadata))


def _blocks_from_text(text: str, metadata: dict[str, Any]) -> tuple[ExtractedBlock, ...]:
    """Recover typed blocks from the stable headings emitted by format extractors."""
    lines = text.splitlines()
    blocks: list[ExtractedBlock] = []
    narrative: list[str] = []
    heading: str | None = None
    table_index = 0

    def flush_narrative() -> None:
        nonlocal narrative
        value = "\n".join(narrative).strip()
        if value:
            blocks.append(
                ExtractedBlock(
                    value,
                    "text",
                    {**metadata, "content_kind": "text", "heading": heading},
                )
            )
        narrative = []

    index = 0
    while index < len(lines):
        line = lines[index]
        if line.lstrip().startswith("|"):
            flush_narrative()
            table_lines: list[str] = []
            while index < len(lines) and lines[index].lstrip().startswith("|"):
                table_lines.append(lines[index])
                index += 1
            table_index += 1
            table_text = "\n".join(table_lines).strip()
            if table_text:
                prefix = f"{heading}\n\n" if heading else ""
                blocks.append(
                    ExtractedBlock(
                        prefix + table_text,
                        "table",
                        {
                            **metadata,
                            "content_kind": "table",
                            "table_index": table_index,
                            "heading": heading,
                        },
                    )
                )
            continue
        if line.lstrip().startswith("#") and line.lstrip().split(" ", 1)[0].strip("#") == "":
            flush_narrative()
            heading = line.strip()
        else:
            narrative.append(line)
        index += 1
    flush_narrative()
    return tuple(blocks)


def _extract_pdf(path: Path, data: bytes) -> ExtractedDocument:
    try:
        import pdfplumber
    except ImportError as exc:
        raise DocumentExtractionError(
            "PDF extraction requires the documents extra: pip install \"recall-rag[documents]\""
        ) from exc

    sections: list[str] = []
    table_count = 0
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page_number, page in enumerate(pdf.pages, start=1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    sections.append(f"## Page {page_number}\n\n{page_text.strip()}")
                for table_number, rows in enumerate(page.extract_tables() or (), start=1):
                    markdown = _rows_to_markdown(rows)
                    if markdown:
                        table_count += 1
                        sections.append(
                            f"### Table {table_count} (page {page_number}, table {table_number})\n\n"
                            f"{markdown}"
                        )
    except Exception as exc:
        raise DocumentExtractionError(f"could not extract PDF {path.name}: {type(exc).__name__}") from exc
    return _result("\n\n".join(sections), "pdf", tables=table_count)


def _extract_docx(data: bytes) -> ExtractedDocument:
    try:
        from docx import Document
    except ImportError as exc:
        raise DocumentExtractionError(
            "DOCX extraction requires the documents extra: pip install \"recall-rag[documents]\""
        ) from exc

    try:
        document = Document(io.BytesIO(data))
        sections = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        tables = 0
        for index, table in enumerate(document.tables, start=1):
            markdown = _rows_to_markdown(
                [[cell.text for cell in row.cells] for row in table.rows]
            )
            if markdown:
                tables += 1
                sections.append(f"### Table {index}\n\n{markdown}")
    except Exception as exc:
        raise DocumentExtractionError(f"could not extract DOCX: {type(exc).__name__}") from exc
    return _result("\n\n".join(sections), "docx", tables=tables)


def _extract_spreadsheet(suffix: str, data: bytes) -> ExtractedDocument:
    if suffix in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise DocumentExtractionError(
                "XLSX extraction requires the documents extra: pip install \"recall-rag[documents]\""
            ) from exc
        try:
            workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=False)
            sections: list[str] = []
            tables = 0
            for sheet in workbook.worksheets:
                rows = _bounded_rows(sheet.iter_rows(values_only=True), sheet.max_row, sheet.max_column)
                markdown = _rows_to_markdown(rows)
                if markdown:
                    tables += 1
                    sections.append(f"## Sheet: {sheet.title}\n\n{markdown}")
            workbook.close()
        except Exception as exc:
            raise DocumentExtractionError(f"could not extract XLSX: {type(exc).__name__}") from exc
        return _result("\n\n".join(sections), suffix.removeprefix("."), tables=tables)

    try:
        import xlrd
    except ImportError as exc:
        raise DocumentExtractionError(
            "XLS extraction requires the documents extra: pip install \"recall-rag[documents]\""
        ) from exc
    try:
        workbook = xlrd.open_workbook(file_contents=data, on_demand=True)
        sections = []
        tables = 0
        for sheet in workbook.sheets():
            sheet_rows = (sheet.row_values(row) for row in range(min(sheet.nrows, MAX_TABLE_ROWS)))
            markdown = _rows_to_markdown(sheet_rows)
            if markdown:
                tables += 1
                sections.append(f"## Sheet: {sheet.name}\n\n{markdown}")
        workbook.release_resources()
    except Exception as exc:
        raise DocumentExtractionError(f"could not extract XLS: {type(exc).__name__}") from exc
    return _result("\n\n".join(sections), "xls", tables=tables)


def _extract_pptx(data: bytes) -> ExtractedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentExtractionError(
            "PPTX extraction requires the documents extra: pip install \"recall-rag[documents]\""
        ) from exc
    try:
        presentation = Presentation(io.BytesIO(data))
        sections: list[str] = []
        tables = 0
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    markdown = _rows_to_markdown(
                        [[cell.text for cell in row.cells] for row in shape.table.rows]
                    )
                    if markdown:
                        tables += 1
                        slide_parts.append(markdown)
                elif getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if text:
                        slide_parts.append(text)
            if slide_parts:
                sections.append(f"## Slide {slide_number}\n\n" + "\n\n".join(slide_parts))
    except Exception as exc:
        raise DocumentExtractionError(f"could not extract PPTX: {type(exc).__name__}") from exc
    return _result("\n\n".join(sections), "pptx", tables=tables)


def _extract_delimited(suffix: str, data: bytes) -> ExtractedDocument:
    delimiter = "\t" if suffix == ".tsv" else ","
    try:
        rows = csv.reader(io.StringIO(_decode_text(data)), delimiter=delimiter)
        markdown = _rows_to_markdown(rows)
    except csv.Error as exc:
        raise DocumentExtractionError(f"could not parse {suffix[1:].upper()}: {exc}") from exc
    return _result(markdown, suffix.removeprefix("."), tables=1)


def _extract_html(data: bytes) -> ExtractedDocument:
    text = _decode_text(data)
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        parser = _VisibleTextParser()
        parser.feed(text)
        return _result(parser.text, "html")
    soup = BeautifulSoup(text, "html.parser")
    sections = [value.strip() for value in soup.stripped_strings]
    tables = 0
    for index, table in enumerate(soup.find_all("table"), start=1):
        markdown = _rows_to_markdown(
            [[cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])] for row in table.find_all("tr")]
        )
        if markdown:
            tables += 1
            sections.append(f"Table {index}\n{markdown}")
    return _result("\n".join(sections), "html", tables=tables)


def _extract_email(data: bytes) -> ExtractedDocument:
    try:
        message = BytesParser(policy=policy.default).parsebytes(data)
        sections = [f"Subject: {message.get('subject', '')}", f"From: {message.get('from', '')}"]
        if message.is_multipart():
            bodies = [part.get_content() for part in message.walk() if part.get_content_type() == "text/plain"]
            sections.extend(str(body) for body in bodies)
        else:
            sections.append(str(message.get_content()))
    except Exception as exc:
        raise DocumentExtractionError(f"could not extract EML: {type(exc).__name__}") from exc
    return _result("\n\n".join(sections), "eml")


def _extract_msg(data: bytes) -> ExtractedDocument:
    try:
        from oxmsg import Message  # type: ignore[import-not-found]
    except ImportError as exc:
        raise DocumentExtractionError(
            "MSG extraction requires the documents extra: pip install \"recall-rag[documents]\""
        ) from exc

    try:
        message = Message.load(data)
        sections = [
            f"Subject: {message.subject}",
            f"From: {message.sender or ''}",
            message.body or "",
        ]
    except Exception as exc:
        raise DocumentExtractionError(f"could not extract MSG: {type(exc).__name__}") from exc
    return _result("\n\n".join(sections), "msg")


def _extract_rtf(data: bytes) -> ExtractedDocument:
    text = _decode_text(data)
    text = re.sub(r"\\'[0-9a-fA-F]{2}", "", text)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)
    text = text.replace("{", "").replace("}", "")
    return _result(text, "rtf")


def _extract_epub(data: bytes) -> ExtractedDocument:
    """Extract EPUB spine content without depending on LibreOffice's optional ebook filters."""
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            names = set(archive.namelist())
            if len(names) > MAX_EPUB_MEMBERS:
                raise ValueError("EPUB has too many archive members")
            total_bytes = 0

            def read_member(name: str) -> bytes:
                nonlocal total_bytes
                info = archive.getinfo(name)
                if info.file_size > MAX_EPUB_MEMBER_BYTES:
                    raise ValueError("EPUB archive member is too large")
                if total_bytes + info.file_size > MAX_EPUB_TOTAL_BYTES:
                    raise ValueError("EPUB archive is too large")
                payload = archive.read(name)
                total_bytes += len(payload)
                if total_bytes > MAX_EPUB_TOTAL_BYTES:
                    raise ValueError("EPUB archive is too large")
                return payload

            container = ElementTree.fromstring(read_member("META-INF/container.xml"))
            rootfile = container.find(
                ".//{urn:oasis:names:tc:opendocument:xmlns:container}rootfile"
            )
            opf_name = rootfile.get("full-path", "") if rootfile is not None else ""
            if not opf_name or opf_name not in names:
                raise ValueError("EPUB container has no readable package file")
            package = ElementTree.fromstring(read_member(opf_name))
            opf_ns = "http://www.idpf.org/2007/opf"
            manifest: dict[str, str] = {}
            for item in package.findall(f"{{{opf_ns}}}manifest/{{{opf_ns}}}item"):
                item_id = item.get("id")
                href = item.get("href")
                if item_id and href:
                    manifest[item_id] = href
            spine = package.find(f"{{{opf_ns}}}spine")
            idrefs = (
                [item.get("idref") for item in spine.findall(f"{{{opf_ns}}}itemref")]
                if spine is not None
                else []
            )
            opf_dir = posixpath.dirname(opf_name)
            sections: list[str] = []
            for idref in idrefs:
                href = manifest.get(idref or "")
                if not href:
                    continue
                target = posixpath.normpath(
                    posixpath.join(opf_dir, urllib.parse.unquote(href.split("#", 1)[0]))
                )
                if target not in names or target.startswith("../"):
                    continue
                parser = _VisibleTextParser()
                parser.feed(_decode_text(read_member(target)))
                if parser.text.strip():
                    sections.append(parser.text)
    except (KeyError, ValueError, ElementTree.ParseError, zipfile.BadZipFile) as exc:
        raise DocumentExtractionError(f"could not extract EPUB: {exc}") from exc
    return _result("\n\n".join(sections), "epub")


# A small pool of reusable LibreOffice user profiles, one process-wide, checked out per conversion.
#
# Every call used to build its own `-env:UserInstallation` inside its own temporary directory, so
# every legacy-format extraction paid a cold profile bootstrap. Reusing profiles is worth real time:
# measured 2026-08-18, Windows 11, LibreOffice 26.2.5.2, five serial extractions fell from a median
# 29.76s to 17.94s, because only the first call still pays the bootstrap and the rest drop from
# about 6s to about 3s each. Full record and the falsified prediction that shaped this design:
# `docs/preregistrations/2026-08-18-libreoffice-profile-reuse.md`.
#
# **A pool rather than one shared profile, because one profile has to be locked and the lock costs
# more than the bootstrap it saves.** Two `soffice --convert-to` processes started at the same
# moment against ONE profile do not both succeed: exactly one converts, and the other exits 1
# having written no output file and nothing at all on stderr, cold profile and warm alike. So a
# single shared profile forces every extraction in the process to serialise, and measured at four
# way concurrency that lost outright, 12.61s and 13.16s serialised against 7.10s and 9.25s for the
# old always-cold code. Distinct profiles used concurrently are fine (zero errors in that same
# measurement), so the pool keeps the reuse and keeps the parallelism.
#
# Profiles are per PROCESS rather than at a fixed path under the temp directory, deliberately: a
# lock cannot span processes, but a distinct path does not need to, so two concurrent `recall index`
# processes cannot collide at all. The price is one bootstrap per process, not one per machine.
#
# Re-measure: `python scripts/bench_libreoffice_profile.py probe --fixtures DIR [--warm]` for the
# collision, and `serial` / `threads` for the two timings.
_DEFAULT_PROFILE_POOL = 4
_pool_lock = threading.Lock()
_profile_pool: queue.LifoQueue[str | None] | None = None
_live_profiles: set[str] = set()
_cleanup_registered = False


def _shared_profile_disabled() -> bool:
    """Escape hatch back to a per-call profile, for a deployment that meets a stuck pooled one."""
    setting = os.environ.get("RECALL_LIBREOFFICE_SHARED_PROFILE", "1").strip().lower()
    return setting in {"0", "false", "no", "off"}


def _profile_pool_size() -> int:
    """How many extractions may run at once. Beyond core count there is nothing left to overlap."""
    raw = os.environ.get("RECALL_LIBREOFFICE_PROFILES", "").strip()
    try:
        return max(1, int(raw))
    except ValueError:
        return max(1, min(_DEFAULT_PROFILE_POOL, os.cpu_count() or 1))


def _profiles() -> queue.LifoQueue[str | None]:
    global _profile_pool
    with _pool_lock:
        if _profile_pool is None:
            # LIFO, not FIFO, and that is the whole trick for the serial case: a returned profile
            # goes back on top, so the next call reuses it instead of working through the remaining
            # empty slots and paying a bootstrap for each. A FIFO pool of four would make a serial
            # workload build all four profiles before it reused any of them.
            pool: queue.LifoQueue[str | None] = queue.LifoQueue()
            for _ in range(_profile_pool_size()):
                pool.put(None)
            global _cleanup_registered
            if not _cleanup_registered:
                atexit.register(_discard_all_profiles)
                _cleanup_registered = True
            _profile_pool = pool
        return _profile_pool


def _discard_all_profiles() -> None:
    with _pool_lock:
        stale = list(_live_profiles)
        _live_profiles.clear()
    for directory in stale:
        shutil.rmtree(directory, ignore_errors=True)


@contextlib.contextmanager
def _libreoffice_profile(scratch: Path) -> Iterator[Path]:
    """Check a user profile out of the pool, blocking while every profile is in use."""
    if _shared_profile_disabled():
        profile = scratch / "profile"
        profile.mkdir()
        yield profile
        return
    pool = _profiles()
    slot = pool.get()
    if slot is None:
        slot = tempfile.mkdtemp(prefix="recall-office-profile-")
        with _pool_lock:
            _live_profiles.add(slot)
    try:
        yield Path(slot)
    except BaseException:
        # A `soffice` that died mid-run can leave a lock behind that poisons every later launch
        # against the same profile, which a per-call profile could never suffer from. Throwing the
        # profile away costs one bootstrap; keeping it costs every call that reuses the slot after.
        #
        # `BaseException`, not `Exception`, and the difference is not pedantry: the slot MUST go
        # back. A Ctrl-C through here would otherwise retire one slot permanently, and enough of
        # them would empty the pool and leave the next `pool.get()` blocking for ever.
        with _pool_lock:
            _live_profiles.discard(slot)
        shutil.rmtree(slot, ignore_errors=True)
        pool.put(None)
        raise
    else:
        pool.put(slot)


def _run_libreoffice(
    executable: str,
    path: Path,
    scratch: Path,
    source: Path,
    convert_to: str,
    verb: str,
) -> None:
    """Run one headless conversion, translating both failure modes into a typed error."""
    try:
        with _libreoffice_profile(scratch) as profile:
            subprocess.run(
                [
                    executable,
                    f"-env:UserInstallation={profile.as_uri()}",
                    "--headless",
                    "--convert-to",
                    convert_to,
                    "--outdir",
                    str(scratch),
                    str(source),
                ],
                check=True,
                capture_output=True,
                text=True,
                # LibreOffice reports in the system locale, so this is the call site most likely to
                # emit a byte the platform codec cannot read. See `recall/desktop/runtime.py` for
                # what that costs: rc=0, `stdout=None`, and no exception.
                encoding="utf-8",
                errors="replace",
                timeout=120,
            )
    except FileNotFoundError as exc:
        raise DocumentExtractionError(
            f"{path.suffix.upper()} needs LibreOffice for extraction; install LibreOffice "
            "or use the modern Office format"
        ) from exc
    except subprocess.SubprocessError as exc:
        raise DocumentExtractionError(
            f"LibreOffice could not {verb} {path.name}: {type(exc).__name__}"
        ) from exc


def _extract_with_libreoffice(path: Path, data: bytes) -> ExtractedDocument:
    executable = _libreoffice_executable()
    if executable is None:
        raise DocumentExtractionError(
            f"{path.suffix.upper()} needs LibreOffice for extraction; install LibreOffice or use "
            "the modern Office format"
        )
    # The scratch directory stays per call even though the profile is now shared: two concurrent
    # extractions must not race on the source file or on the converted output name.
    with tempfile.TemporaryDirectory(prefix="recall-office-") as directory:
        scratch = Path(directory)
        source = scratch / path.name
        source.write_bytes(data)
        if path.suffix.lower() in {".ppt", ".odp"}:
            _run_libreoffice(executable, path, scratch, source, "pptx", "convert")
            converted_pptx = source.with_suffix(".pptx")
            if not converted_pptx.exists():
                raise DocumentExtractionError(f"LibreOffice produced no PPTX for {path.name}")
            return _extract_pptx(converted_pptx.read_bytes())
        output_suffix = ".csv" if path.suffix.lower() == ".ods" else ".txt"
        output_filter = "csv:Text - txt - csv (StarCalc)" if output_suffix == ".csv" else "txt:Text"
        _run_libreoffice(executable, path, scratch, source, output_filter, "extract")
        converted = source.with_suffix(output_suffix)
        if not converted.exists():
            raise DocumentExtractionError(f"LibreOffice produced no text for {path.name}")
        if output_suffix == ".csv":
            try:
                markdown = _rows_to_markdown(csv.reader(io.StringIO(_decode_text(converted.read_bytes()))))
            except csv.Error as exc:
                raise DocumentExtractionError(f"could not parse ODS {path.name}: {exc}") from exc
            return _result(markdown, "ods", tables=1)
        return _result(_decode_text(converted.read_bytes()), path.suffix.removeprefix("."))


def _libreoffice_executable() -> str | None:
    """Find the LibreOffice CLI without requiring users to edit PATH on desktop installs."""
    candidates = [
        os.environ.get("RECALL_LIBREOFFICE"),
        shutil.which("soffice"),
        shutil.which("libreoffice"),
    ]
    if os.name == "nt":
        for variable in ("PROGRAMFILES", "PROGRAMFILES(X86)"):
            root = os.environ.get(variable)
            if root:
                candidates.append(str(Path(root) / "LibreOffice" / "program" / "soffice.exe"))
    for candidate in candidates:
        if candidate and Path(candidate).is_file():
            return candidate
    return None


def _bounded_rows(rows: Iterable[Iterable[Any]], max_row: int | None, max_column: int | None) -> list[list[Any]]:
    row_limit = min(max_row or MAX_TABLE_ROWS, MAX_TABLE_ROWS)
    column_limit = min(max_column or MAX_TABLE_COLUMNS, MAX_TABLE_COLUMNS)
    result: list[list[Any]] = []
    for row_index, row in enumerate(rows):
        if row_index >= row_limit or len(result) * column_limit >= MAX_TABLE_CELLS:
            break
        result.append(list(row)[:column_limit])
    return result


def _rows_to_markdown(rows: Iterable[Iterable[Any]]) -> str:
    normalized = [[_cell(value) for value in row] for row in rows]
    normalized = [row for row in normalized if any(value for value in row)]
    if not normalized:
        return ""
    width = min(max(len(row) for row in normalized), MAX_TABLE_COLUMNS)
    normalized = [(row + [""] * width)[:width] for row in normalized]
    header = normalized[0]
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join("---" for _ in header) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in normalized[1:])
    return "\n".join(lines)


def _cell(value: Any) -> str:
    if value is None:
        return ""
    text = str(value).replace("\r", " ").replace("\n", " ").strip()
    return text.replace("|", "\\|")


def chunk_extracted_document(
    document: ExtractedDocument,
    *,
    max_chars: int = 800,
    overlap: int = 80,
) -> list[tuple[str, dict[str, Any]]]:
    """Chunk extracted material while preserving table headers and provenance.

    Narrative blocks use paragraph packing. Table blocks use repeated headers and row groups,
    which keeps a row meaningful when it is retrieved without its neighbours.
    """
    if max_chars < 1:
        raise ValueError("max_chars must be >= 1")
    chunks: list[tuple[str, dict[str, Any]]] = []
    for block in document.blocks or (
        ExtractedBlock(document.text, "text", dict(document.metadata)),
    ):
        if block.kind == "table":
            chunks.extend(_chunk_table_block(block, max_chars=max_chars))
            continue
        pieces = _chunk_text_fallback(block.text, max_chars=max_chars, overlap=overlap)
        for piece in pieces:
            chunks.append((piece, {**block.metadata, "content_kind": "text"}))
    return chunks


def _chunk_table_block(
    block: ExtractedBlock, *, max_chars: int
) -> list[tuple[str, dict[str, Any]]]:
    lines = [line.strip() for line in block.text.splitlines() if line.strip()]
    pipe_lines = [line for line in lines if line.startswith("|")]
    if len(pipe_lines) < 2:
        return [(block.text[:max_chars], {**block.metadata, "content_kind": "table"})]
    prefix = [line for line in lines if not line.startswith("|")]
    header = pipe_lines[:2]
    rows = pipe_lines[2:]
    prefix_text = "\n\n".join(prefix)
    header_text = "\n".join(header)
    base = f"{prefix_text}\n\n{header_text}" if prefix_text else header_text
    table_metadata = {"table_headers": _table_header_values(header[0])}
    result: list[tuple[str, dict[str, Any]]] = []
    current: list[str] = []
    current_length = 0
    start_row = 0

    def emit(end_row: int) -> None:
        if not current:
            return
        text = base + "\n" + "\n".join(current)
        metadata = {
            **block.metadata,
            **table_metadata,
            "numeric_values": _table_numeric_values(current),
            "content_kind": "table",
            "row_start": start_row,
            "row_end": end_row,
        }
        if len(text) <= max_chars:
            result.append((text, metadata))
            return
        if len(base) >= max_chars:
            for piece in _chunk_text_fallback(text, max_chars=max_chars, overlap=0):
                result.append((piece, metadata))
            return
        budget = max(1, max_chars - len(base) - 1)
        for piece in _chunk_text_fallback("\n".join(current), max_chars=budget, overlap=0):
            result.append((base + "\n" + piece, metadata))

    for row_index, row in enumerate(rows):
        candidate_length = len(base) + 1 + current_length + (1 if current else 0) + len(row)
        if current and candidate_length > max_chars:
            emit(row_index - 1)
            current = [row]
            current_length = len(row)
            start_row = row_index
        else:
            current.append(row)
            current_length += len(row) + (1 if current_length else 0)
    emit(len(rows) - 1)
    return result or [
        (base, {**block.metadata, **table_metadata, "content_kind": "table"})
    ]


def _table_header_values(line: str) -> list[str]:
    return [value.strip() for value in line.strip("|").split("|") if value.strip()]


def _table_numeric_values(rows: list[str]) -> list[str]:
    values: list[str] = []
    for row in rows:
        for cell in row.strip("|").split("|"):
            for match in re.findall(r"(?<![\w.])[+-]?\d+(?:[.,]\d+)?%?(?![\w.])", cell):
                normalized = match.replace(",", ".")
                if normalized not in values:
                    values.append(normalized)
    return values


def _chunk_text_fallback(text: str, *, max_chars: int, overlap: int) -> list[str]:
    paragraphs = [part.strip() for part in text.split("\n\n") if part.strip()]
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        if len(paragraph) > max_chars:
            if current:
                chunks.append(current)
                current = ""
            start = 0
            while start < len(paragraph):
                end = min(start + max_chars, len(paragraph))
                if end < len(paragraph):
                    boundary = max(paragraph.rfind(" ", start, end), paragraph.rfind("\n", start, end))
                    if boundary > start:
                        end = boundary
                piece = paragraph[start:end].strip()
                if piece:
                    chunks.append(piece)
                if end >= len(paragraph):
                    break
                start = max(end - min(overlap, max_chars // 4), start + 1)
            continue
        candidate = f"{current}\n\n{paragraph}" if current else paragraph
        if current and len(candidate) > max_chars:
            chunks.append(current)
            current = paragraph
        else:
            current = candidate
    if current:
        chunks.append(current)
    return chunks


class _VisibleTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        value = data.strip()
        if value:
            self.parts.append(value)

    @property
    def text(self) -> str:
        return "\n".join(self.parts)
