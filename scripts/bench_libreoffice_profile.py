"""Measure the LibreOffice profile reuse questions pre-registered on 2026-08-18.

Record: ``docs/preregistrations/2026-08-18-libreoffice-profile-reuse.md``.

Four subcommands, one per question in that record:

``fixtures``    build the five legacy-format files once, so the timed sections measure extraction
                and never fixture conversion.
``serial``      Q1. Time five ``extract_document`` calls, one per LibreOffice format, in one process.
``probe``       Q2. Launch two ``soffice --convert-to`` processes against one shared
                ``-env:UserInstallation`` at the same moment and report what actually happens.
``threads``     Q3. Four concurrent ``extract_document`` calls, wall clock.

Every subcommand prints one JSON object on stdout, so a driver can collect runs without parsing
prose. Arms are selected with ``RECALL_LIBREOFFICE_SHARED_PROFILE``, which is the same switch a
deployment would use, rather than by editing the module under measurement.
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
import tempfile
import threading
import time
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from recall.extraction import _libreoffice_executable, extract_document  # noqa: E402

# Source fixture against target legacy format, in the order the pre-registered series reports them.
CASES: list[tuple[str, str]] = [
    ("source.docx", "doc"),
    ("source.docx", "odt"),
    ("source.xlsx", "ods"),
    ("source.pptx", "ppt"),
    ("source.pptx", "odp"),
]

# The marker each fixture carries, so a conversion that silently produced an empty document is
# caught here rather than scored as a fast extraction.
EXPECTED: dict[str, str] = {
    "doc": "LEGACY-DOC-TEST",
    "odt": "LEGACY-DOC-TEST",
    "ods": "123.45",
    "ppt": "LEGACY-PPT-TEST",
    "odp": "LEGACY-PPT-TEST",
}


def _libreoffice_version(executable: str) -> str:
    """`soffice --version` prints nothing on Windows, so read the binary's own VersionInfo."""
    if sys.platform != "win32":
        result = subprocess.run([executable, "--version"], capture_output=True, text=True)
        return result.stdout.strip() or "unknown"
    result = subprocess.run(
        [
            "powershell",
            "-NoProfile",
            "-Command",
            f"(Get-Item '{executable}').VersionInfo.ProductVersion",
        ],
        capture_output=True,
        text=True,
    )
    return result.stdout.strip() or "unknown"


def build_fixtures(outdir: Path) -> dict[str, object]:
    """Convert modern fixtures down to the five legacy formats, with a throwaway profile."""
    import docx
    import openpyxl
    import pptx

    executable = _libreoffice_executable()
    if executable is None:
        raise SystemExit("LibreOffice is not installed")

    outdir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="recall-bench-setup-") as scratch:
        sources = Path(scratch) / "sources"
        profile = Path(scratch) / "profile"
        sources.mkdir()
        profile.mkdir()

        document = docx.Document()
        document.add_paragraph("LEGACY-DOC-TEST")
        document.save(sources / "source.docx")

        workbook = openpyxl.Workbook()
        workbook.active.append(["Region", "Revenue"])
        workbook.active.append(["EU", 123.45])
        workbook.save(sources / "source.xlsx")

        presentation = pptx.Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "LEGACY-PPT-TEST"
        presentation.save(sources / "source.pptx")

        for source_name, target in CASES:
            subprocess.run(
                [
                    executable,
                    f"-env:UserInstallation={profile.as_uri()}",
                    "--headless",
                    "--convert-to",
                    target,
                    "--outdir",
                    str(outdir),
                    str(sources / source_name),
                ],
                check=True,
                capture_output=True,
                text=True,
                timeout=180,
            )
            produced = outdir / f"{Path(source_name).stem}.{target}"
            if not produced.exists():
                raise SystemExit(f"fixture build produced no {target}")

    return {
        "fixtures": sorted(p.name for p in outdir.iterdir()),
        "libreoffice_version": _libreoffice_version(executable),
    }


def _fixture_path(fixtures: Path, source_name: str, target: str) -> Path:
    return fixtures / f"{Path(source_name).stem}.{target}"


def run_serial(fixtures: Path, arm: str) -> dict[str, object]:
    """Q1. Five extractions in order, timed individually."""
    per_call: list[float] = []
    texts: dict[str, str] = {}
    for source_name, target in CASES:
        path = _fixture_path(fixtures, source_name, target)
        data = path.read_bytes()
        started = time.perf_counter()
        document = extract_document(path, data)
        per_call.append(round(time.perf_counter() - started, 2))
        texts[target] = document.text
        if EXPECTED[target] not in document.text:
            raise SystemExit(f"{target}: expected marker missing, the timing is meaningless")
    return {
        "question": "Q1",
        "arm": arm,
        "per_call_seconds": per_call,
        "total_extraction_seconds": round(sum(per_call), 2),
        "n": len(per_call),
        # Q4 rides along: the caller compares these across arms.
        "text_sha": {k: str(hash(v)) for k, v in texts.items()},
        "texts": texts,
    }


def run_probe(fixtures: Path, warm: bool) -> dict[str, object]:
    """Q2. Two soffice conversions sharing one profile, started at the same moment."""
    executable = _libreoffice_executable()
    if executable is None:
        raise SystemExit("LibreOffice is not installed")

    with tempfile.TemporaryDirectory(prefix="recall-bench-probe-") as scratch:
        root = Path(scratch)
        profile = root / "shared-profile"
        profile.mkdir()

        # Two independent inputs and two independent output directories, so the only thing the two
        # processes share is the user profile. Anything that goes wrong is attributable to that.
        pairs = []
        for index, (source_name, target) in enumerate(CASES[:2]):
            source = root / f"input{index}.{target}"
            source.write_bytes(_fixture_path(fixtures, source_name, target).read_bytes())
            outdir = root / f"out{index}"
            outdir.mkdir()
            pairs.append((source, outdir))

        def command(source: Path, outdir: Path) -> list[str]:
            return [
                executable,
                f"-env:UserInstallation={profile.as_uri()}",
                "--headless",
                "--convert-to",
                "txt:Text",
                "--outdir",
                str(outdir),
                str(source),
            ]

        warmed_seconds = None
        if warm:
            # Bootstrap the profile first, so the concurrent pair meets an already-built profile
            # rather than racing to create one. These are different failure modes.
            started = time.perf_counter()
            subprocess.run(command(*pairs[0]), capture_output=True, text=True, timeout=180)
            warmed_seconds = round(time.perf_counter() - started, 2)
            for _, outdir in pairs:
                for stale in outdir.iterdir():
                    stale.unlink()

        started = time.perf_counter()
        processes = [
            subprocess.Popen(command(source, outdir), stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            for source, outdir in pairs
        ]
        results = []
        for index, process in enumerate(processes):
            stdout, stderr = process.communicate(timeout=180)
            source, outdir = pairs[index]
            produced = outdir / f"{source.stem}.txt"
            results.append(
                {
                    "index": index,
                    "returncode": process.returncode,
                    "stdout": stdout.decode("utf-8", "replace").strip()[:400],
                    "stderr": stderr.decode("utf-8", "replace").strip()[:400],
                    "output_exists": produced.exists(),
                    "output_bytes": produced.stat().st_size if produced.exists() else 0,
                }
            )
        elapsed = round(time.perf_counter() - started, 2)

    both_ok = all(r["returncode"] == 0 and r["output_exists"] and r["output_bytes"] > 0 for r in results)
    return {
        "question": "Q2",
        "profile": "warm" if warm else "cold",
        "warm_bootstrap_seconds": warmed_seconds,
        "concurrent_seconds": elapsed,
        "processes": results,
        "both_succeeded": both_ok,
    }


def run_threads(fixtures: Path, arm: str, workers: int) -> dict[str, object]:
    """Q3. Four concurrent extractions, wall clock."""
    cases = [CASES[index % len(CASES)] for index in range(workers)]
    barrier = threading.Barrier(workers)

    def one(case: tuple[str, str]) -> dict[str, object]:
        source_name, target = case
        path = _fixture_path(fixtures, source_name, target)
        data = path.read_bytes()
        barrier.wait()  # start all workers at the same moment, not as the pool schedules them
        started = time.perf_counter()
        try:
            document = extract_document(path, data)
        except Exception as exc:  # a collision is a result here, not a crash
            return {"target": target, "error": f"{type(exc).__name__}: {exc}"}
        return {
            "target": target,
            "seconds": round(time.perf_counter() - started, 2),
            "marker_found": EXPECTED[target] in document.text,
        }

    started = time.perf_counter()
    with ThreadPoolExecutor(max_workers=workers) as pool:
        outcomes = list(pool.map(one, cases))
    elapsed = round(time.perf_counter() - started, 2)
    return {
        "question": "Q3",
        "arm": arm,
        "workers": workers,
        "wall_clock_seconds": elapsed,
        "outcomes": outcomes,
        "errors": sum(1 for o in outcomes if "error" in o),
    }


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    sub = parser.add_subparsers(dest="command", required=True)

    fixtures_parser = sub.add_parser("fixtures")
    fixtures_parser.add_argument("--outdir", type=Path, required=True)

    serial_parser = sub.add_parser("serial")
    serial_parser.add_argument("--fixtures", type=Path, required=True)
    serial_parser.add_argument("--arm", required=True)

    probe_parser = sub.add_parser("probe")
    probe_parser.add_argument("--fixtures", type=Path, required=True)
    probe_parser.add_argument("--warm", action="store_true")

    threads_parser = sub.add_parser("threads")
    threads_parser.add_argument("--fixtures", type=Path, required=True)
    threads_parser.add_argument("--arm", required=True)
    threads_parser.add_argument("--workers", type=int, default=4)

    args = parser.parse_args()
    if args.command == "fixtures":
        payload = build_fixtures(args.outdir)
    elif args.command == "serial":
        payload = run_serial(args.fixtures, args.arm)
    elif args.command == "probe":
        payload = run_probe(args.fixtures, args.warm)
    else:
        payload = run_threads(args.fixtures, args.arm, args.workers)
    print(json.dumps(payload, indent=2))


if __name__ == "__main__":
    main()
